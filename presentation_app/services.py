"""
Marp CLIを操作するサービスモジュール

このモジュールは、Marp CLIを使用してMarkdownをプレゼンテーション形式に変換する機能を提供します。
"""

import os
import subprocess
import tempfile
import platform
import shutil
import re
import random
import math
from django.conf import settings
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class MarpService:
    """
    Marp CLI操作のためのサービスクラス
    
    Markdownをプレゼンテーション形式（PPTX）に変換するためのサービスメソッドを提供します。
    """
    
    # 1スライドあたりの文字数の目安（フォントサイズ別）
    CHARS_PER_SLIDE = {
        'large': 100,    # 大きいフォント向け（タイトルスライドなど）
        'medium': 200,   # 中くらいのフォント向け（一般的なコンテンツスライド）
        'small': 300     # 小さいフォント向け（詳細情報スライド）
    }
    
    @staticmethod
    def markdown_to_pptx(markdown_content, theme='default'):
        """
        MarkdownをMarp CLIを使ってPPTXに変換する
        
        Args:
            markdown_content (str): 変換する元のMarkdown内容
            theme (str): 使用するMarpテーマ名（デフォルト: 'default'）
            
        Returns:
            bytes: 生成されたPPTXファイルのバイナリ内容、失敗時はNone
        """
        # 一時ディレクトリを作成
        temp_dir = tempfile.mkdtemp()
        markdown_path = os.path.join(temp_dir, 'presentation.md')
        output_path = os.path.join(temp_dir, 'presentation.pptx')
        
        try:
            # Markdownファイルの先頭にMarpディレクティブを追加
            marp_directives = f"""---
marp: true
theme: {theme}
---

"""
            full_markdown = marp_directives + markdown_content
            
            # Markdownファイルを作成
            with open(markdown_path, 'w', encoding='utf-8') as f:
                f.write(full_markdown)
            
            # Windowsの場合はshellをTrueに設定
            is_windows = platform.system() == 'Windows'
            
            # 3つの方法でMarp CLIの実行を試みる
            pptx_content = None
            
            # 方法1: npxを使用
            if pptx_content is None:
                try:
                    pptx_content = MarpService._try_npx_marp_cli(markdown_path, output_path, is_windows)
                except Exception as e:
                    print(f"npxでのMarp CLI実行に失敗: {e}")
            
            # 方法2: グローバルにインストールされているMarp CLIを使用
            if pptx_content is None:
                try:
                    pptx_content = MarpService._try_global_marp_cli(markdown_path, output_path, is_windows)
                except Exception as e:
                    print(f"グローバルMarp CLI実行に失敗: {e}")
            
            # 方法3: node_modulesから直接実行
            if pptx_content is None:
                try:
                    pptx_content = MarpService._try_local_marp_cli(markdown_path, output_path, is_windows)
                except Exception as e:
                    print(f"ローカルMarp CLI実行に失敗: {e}")
            
            if pptx_content:
                return pptx_content
            else:
                print("すべてのMarp CLI実行方法が失敗しました")
                return None
            
        except Exception as e:
            print(f"予期せぬエラー: {e}")
            return None
        finally:
            # 一時ファイルの削除（clean up）
            try:
                if os.path.exists(markdown_path):
                    os.remove(markdown_path)
                if os.path.exists(output_path):
                    os.remove(output_path)
                if os.path.exists(temp_dir):
                    os.rmdir(temp_dir)
            except Exception as e:
                print(f"一時ファイルの削除中にエラーが発生しました: {e}")
                # エラーが発生しても処理を続行する
    
    @staticmethod
    def _try_npx_marp_cli(markdown_path, output_path, is_windows):
        """npxを使用してMarp CLIを実行する"""
        cmd = [
            'npx',
            '@marp-team/marp-cli',
            markdown_path,
            '--output', output_path,
            '--pptx',
            '--pptx-convert-infix', 'false'
        ]
        
        result = subprocess.run(cmd, check=True, shell=is_windows, capture_output=True)
        print(f"npx Marp CLI実行成功: {result.stdout.decode('utf-8', errors='ignore')}")
        
        if os.path.exists(output_path):
            with open(output_path, 'rb') as f:
                return f.read()
        return None
    
    @staticmethod
    def _try_global_marp_cli(markdown_path, output_path, is_windows):
        """グローバルにインストールされているMarp CLIを使用する"""
        cmd = [
            'marp',
            markdown_path,
            '--output', output_path,
            '--pptx',
            '--pptx-convert-infix', 'false'
        ]
        
        result = subprocess.run(cmd, check=True, shell=is_windows, capture_output=True)
        print(f"グローバルMarp CLI実行成功: {result.stdout.decode('utf-8', errors='ignore')}")
        
        if os.path.exists(output_path):
            with open(output_path, 'rb') as f:
                return f.read()
        return None
    
    @staticmethod
    def _try_local_marp_cli(markdown_path, output_path, is_windows):
        """プロジェクトのnode_modulesからMarp CLIを直接実行する"""
        # プロジェクトルートからの相対パス
        marp_cli_path = os.path.join(settings.BASE_DIR, 'node_modules', '.bin', 'marp')
        
        # Windowsの場合は.cmdファイルを使用
        if is_windows and not os.path.exists(marp_cli_path):
            marp_cli_path = f"{marp_cli_path}.cmd"
        
        # パスが存在するか確認
        if not os.path.exists(marp_cli_path):
            print(f"ローカルのMarp CLIが見つかりません: {marp_cli_path}")
            return None
        
        cmd = [
            marp_cli_path,
            markdown_path,
            '--output', output_path,
            '--pptx',
            '--pptx-convert-infix', 'false'
        ]
        
        result = subprocess.run(cmd, check=True, shell=is_windows, capture_output=True)
        print(f"ローカルMarp CLI実行成功: {result.stdout.decode('utf-8', errors='ignore')}")
        
        if os.path.exists(output_path):
            with open(output_path, 'rb') as f:
                return f.read()
        return None
        
    @staticmethod
    def natural_language_to_pptx(title, content, template_path=None):
        """
        自然言語からPPTXを生成する
        
        自然言語の内容を解析し、プレゼンテーションスライドに変換します。
        テンプレートが指定されている場合は、そのテンプレートを基に生成します。
        
        Args:
            title (str): プレゼンテーションのタイトル
            content (str): 自然言語で記述されたプレゼンテーション内容
            template_path (str, optional): 使用するテンプレートPPTXのパス
            
        Returns:
            bytes: 生成されたPPTXファイルのバイナリ内容、失敗時はNone
        """
        try:
            # 一時ディレクトリを作成
            temp_dir = tempfile.mkdtemp()
            output_path = os.path.join(temp_dir, 'presentation.pptx')
            
            # テンプレートがある場合はそれをコピーし、編集する
            prs = MarpService._create_presentation_with_template(template_path)
            
            # スライド内容を解析し、文字量に基づいて分割する
            slides_content = MarpService._parse_natural_language_to_slides_with_sizing(content)
            
            # テンプレートの情報を取得
            template_layouts = MarpService._analyze_template_layouts(prs)
            max_slides = len(prs.slides)
            
            # 既存のスライドを全て削除
            for _ in range(len(prs.slides)):
                r_id = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(r_id)
                prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
            
            # タイトルスライドを追加
            if template_layouts and 'title' in template_layouts:
                title_layout = template_layouts['title']['layout']
            else:
                title_layout = prs.slide_layouts[0]  # デフォルトのタイトルレイアウト
            
            title_slide = prs.slides.add_slide(title_layout)
            
            # タイトルとサブタイトルを設定
            if hasattr(title_slide, 'shapes') and hasattr(title_slide.shapes, 'title') and title_slide.shapes.title:
                title_slide.shapes.title.text = title
            
            # サブタイトルがある場合は設定
            for shape in title_slide.placeholders:
                if shape.placeholder_format.type == 2:  # 2はサブタイトル
                    shape.text = "作成日: " + MarpService._get_current_date_japanese()
            
            # 内容スライドを追加
            for i, slide_content in enumerate(slides_content):
                # 必要なスライド数が最大数を超える場合、ランダムにレイアウトを選択
                if i >= max_slides - 1:  # -1はタイトルスライド分
                    # コンテンツスライドのレイアウトをランダムに選択
                    content_layouts = [l for k, l in template_layouts.items() if k != 'title']
                    if content_layouts:
                        selected_layout = random.choice(content_layouts)
                        content_slide_layout = selected_layout['layout']
                    else:
                        content_slide_layout = prs.slide_layouts[1]  # デフォルトのコンテンツレイアウト
                else:
                    # テンプレートのレイアウトを順番に使用
                    layout_index = (i + 1) % max_slides  # +1はタイトルスライドをスキップ
                    content_slide_layout = prs.slide_masters[0].slide_layouts[layout_index]
                
                slide = prs.slides.add_slide(content_slide_layout)
                
                # タイトルを設定
                if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title') and slide.shapes.title:
                    slide.shapes.title.text = slide_content['title']
                
                # 内容を設定
                content_placeholder = None
                for shape in slide.placeholders:
                    if shape.placeholder_format.type != 1:  # 1はタイトル
                        content_placeholder = shape
                        break
                
                if content_placeholder:
                    text_frame = content_placeholder.text_frame
                    text_frame.clear()  # 既存のテキストをクリア
                    
                    for i, point in enumerate(slide_content['points']):
                        if i == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        p.text = point
                        p.level = 0  # 最上位レベル
            
            # PPTXファイルを保存
            prs.save(output_path)
            
            # 生成されたPPTXファイルを読み込む
            with open(output_path, 'rb') as f:
                pptx_content = f.read()
            
            return pptx_content
            
        except Exception as e:
            print(f"PPTXの生成中にエラーが発生しました: {e}")
            import traceback
            traceback.print_exc()
            return None
            
        finally:
            # 一時ファイルの削除
            try:
                if os.path.exists(output_path):
                    os.remove(output_path)
                if os.path.exists(temp_dir):
                    os.rmdir(temp_dir)
            except Exception as e:
                print(f"一時ファイルの削除中にエラーが発生しました: {e}")
    
    @staticmethod
    def _create_presentation_with_template(template_path):
        """
        テンプレートをベースにプレゼンテーションを作成
        
        Args:
            template_path (str, optional): 使用するテンプレートPPTXのパス
            
        Returns:
            Presentation: 作成されたプレゼンテーションオブジェクト
        """
        if template_path and os.path.exists(template_path):
            prs = Presentation(template_path)
            print(f"テンプレート '{template_path}' を使用してプレゼンテーションを作成します")
        else:
            prs = Presentation()
            print(f"新規プレゼンテーションを作成します（テンプレートなし）")
        
        return prs
    
    @staticmethod
    def _analyze_template_layouts(presentation):
        """
        テンプレートのレイアウト情報を解析
        
        Args:
            presentation (Presentation): 解析するプレゼンテーション
            
        Returns:
            dict: レイアウト情報
        """
        layouts = {}
        
        # スライドマスターからレイアウトを取得
        for i, slide_layout in enumerate(presentation.slide_masters[0].slide_layouts):
            layout_info = {
                'layout': slide_layout,
                'placeholders': {}
            }
            
            # プレースホルダーを分析
            for placeholder in slide_layout.placeholders:
                ph_type = placeholder.placeholder_format.type
                layout_info['placeholders'][ph_type] = {
                    'idx': placeholder.placeholder_format.idx,
                    'type': ph_type,
                    'name': placeholder.name
                }
            
            # タイトルスライドの判定
            is_title_slide = False
            for ph in layout_info['placeholders'].values():
                if ph['name'] and ('タイトル' in ph['name'] or 'Title' in ph['name']):
                    is_title_slide = True
            
            if is_title_slide and i == 0:
                layouts['title'] = layout_info
            else:
                layouts[f'content_{i}'] = layout_info
        
        return layouts
    
    @staticmethod
    def _parse_natural_language_to_slides_with_sizing(content):
        """
        自然言語のテキストをスライド形式に解析し、文字量に基づいて分割
        
        Args:
            content (str): 自然言語のテキスト
            
        Returns:
            list: スライド情報のリスト（文字量に基づいて分割）
        """
        # まず通常の方法でスライドを解析
        slides = MarpService._parse_natural_language_to_slides(content)
        sized_slides = []
        
        # 各スライドを文字量に基づいて分割
        for slide in slides:
            title = slide['title']
            points = slide['points']
            
            # ポイントの総文字数を計算
            total_chars = sum(len(point) for point in points)
            
            # 文字の大きさに基づいて1スライドあたりの文字数を決定
            if "概要" in title or "まとめ" in title or len(title) <= 10:
                chars_per_slide = MarpService.CHARS_PER_SLIDE['large']
                font_size = 'large'
            elif len(points) > 0 and total_chars / len(points) >= 50:  # 平均文字数が多い場合
                chars_per_slide = MarpService.CHARS_PER_SLIDE['small']
                font_size = 'small'
            else:
                chars_per_slide = MarpService.CHARS_PER_SLIDE['medium']
                font_size = 'medium'
            
            # 必要なスライド数を計算
            num_slides_needed = 1
            if total_chars > 0:
                num_slides_needed = math.ceil(total_chars / chars_per_slide)
            
            if num_slides_needed <= 1:
                # 1スライドで足りる場合はそのまま追加
                sized_slides.append({
                    'title': title,
                    'points': points,
                    'font_size': font_size
                })
            else:
                # 複数スライドに分割する必要がある場合
                points_per_slide = math.ceil(len(points) / num_slides_needed)
                
                for i in range(num_slides_needed):
                    start_idx = i * points_per_slide
                    end_idx = min((i + 1) * points_per_slide, len(points))
                    slide_points = points[start_idx:end_idx]
                    
                    # スライド番号を付与（2枚目以降）
                    if i > 0:
                        slide_title = f"{title} ({i+1}/{num_slides_needed})"
                    else:
                        slide_title = title
                    
                    sized_slides.append({
                        'title': slide_title,
                        'points': slide_points,
                        'font_size': font_size
                    })
        
        return sized_slides
    
    @staticmethod
    def _parse_natural_language_to_slides(content):
        """
        自然言語のテキストをスライド形式に解析する
        
        Args:
            content (str): 自然言語のテキスト
            
        Returns:
            list: スライド情報のリスト
        """
        slides = []
        
        # 段落で分割
        paragraphs = content.split('\n\n')
        
        current_slide = None
        
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
                
            # 新しいスライドの開始（「スライド:」または「テーマ:」で始まる行）
            if para.lower().startswith(('スライド:', 'slide:', 'テーマ:', 'theme:')):
                # 前のスライドがあれば追加
                if current_slide:
                    slides.append(current_slide)
                
                # タイトルを抽出
                title = para.split(':', 1)[1].strip()
                current_slide = {
                    'title': title,
                    'points': []
                }
            # 箇条書きのポイント
            elif para.startswith(('-', '・', '*', '•')):
                if not current_slide:
                    current_slide = {
                        'title': '内容',
                        'points': []
                    }
                point = para[1:].strip()
                current_slide['points'].append(point)
            # 通常のテキスト
            else:
                if not current_slide:
                    # 最初の段落はタイトルとして扱う
                    current_slide = {
                        'title': para,
                        'points': []
                    }
                else:
                    # それ以外は箇条書きとして追加
                    sentences = re.split(r'(?<=[。．.!?])\s*', para)
                    for sentence in sentences:
                        if sentence.strip():
                            current_slide['points'].append(sentence.strip())
        
        # 最後のスライドを追加
        if current_slide:
            slides.append(current_slide)
            
        return slides
    
    @staticmethod
    def _get_current_date_japanese():
        """現在の日付を日本語形式で返す"""
        from datetime import datetime
        now = datetime.now()
        return f"{now.year}年{now.month}月{now.day}日"
    
    @staticmethod
    def edit_template_pptx(new_template_path=None):
        """
        テンプレートPPTXを編集する
        
        既存のテンプレートPPTXを編集し、新しいテンプレートとして保存します。
        新しいテンプレートのパスが指定されていない場合は、既存のテンプレートを上書きします。
        
        Args:
            new_template_path (str, optional): 新しいテンプレートの保存パス
            
        Returns:
            bool: 成功した場合はTrue、失敗した場合はFalse
        """
        template_path = os.path.join(settings.BASE_DIR, 'template.pptx')
        
        if not os.path.exists(template_path):
            print(f"テンプレートファイルが見つかりません: {template_path}")
            return False
        
        try:
            # テンプレートを読み込む
            prs = Presentation(template_path)
            
            # 保存先パスを決定
            if not new_template_path:
                # バックアップを作成
                backup_path = os.path.join(settings.BASE_DIR, 'template_backup.pptx')
                shutil.copy2(template_path, backup_path)
                save_path = template_path
            else:
                save_path = new_template_path
            
            # テンプレートを保存
            prs.save(save_path)
            print(f"テンプレートを保存しました: {save_path}")
            
            return True
            
        except Exception as e:
            print(f"テンプレートの編集中にエラーが発生しました: {e}")
            return False
        
    @staticmethod
    def get_template_info():
        """
        テンプレートPPTXのスライド情報を取得する
        
        Returns:
            dict: テンプレート情報
        """
        template_path = os.path.join(settings.BASE_DIR, 'template.pptx')
        
        if not os.path.exists(template_path):
            return {
                'exists': False,
                'message': 'テンプレートファイルが見つかりません'
            }
        
        try:
            prs = Presentation(template_path)
            
            slides_info = []
            for i, slide in enumerate(prs.slides):
                slide_info = {
                    'index': i,
                    'layout_name': f"レイアウト {i+1}",
                }
                
                # タイトルがある場合は取得
                if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title') and slide.shapes.title:
                    slide_info['sample_title'] = slide.shapes.title.text
                
                slides_info.append(slide_info)
            
            return {
                'exists': True,
                'path': template_path,
                'slide_count': len(prs.slides),
                'slides': slides_info
            }
            
        except Exception as e:
            return {
                'exists': False,
                'message': f'テンプレートの読み込み中にエラーが発生しました: {str(e)}'
            } 